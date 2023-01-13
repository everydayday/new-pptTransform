from pptx import Presentation 
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR



def make_slide(lyric_str):
    
    
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # slide 마다 가사 리스트
    slide_str = get_slide_str(lyric_str)
    
    # slide 갯수 구하기 (올림)
    global slide_num
    slide_num = len(slide_str)
    
    # slide 만들기
    for i in range(0, slide_num):
        blank_slide_layout = prs.slide_layouts[6]  # 슬라이드 종류 선택
        slide = prs.slides.add_slide(blank_slide_layout)  # 슬라이드 추가

        # 위치, 가로/세로 길이 - 텍스트 상자
        left = Cm(7.62)
        top = Cm(0)
        width = Cm(25.4)
        height = Cm(19)
        
        
        # 텍스트 상자 넣기
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.text = slide_str[i]
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # TOP 부분만 나와서 MIDDLE 찍었는데 이게 되네...
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].align = PP_ALIGN.CENTER
        
        for paragraph in tf.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
        
        
    prs.save('lyricsPPT.pptx')
    
    
    
    
    
# slide_str 만들기 / index : 슬라이드마다 들어가는 가사 두줄
def get_slide_str(lyric_str):
    my_list = lyric_str.split('\n')
    slide_list = []     # slide 마다 들어가는 lyric list
    idx_count = 0

    for idx, value in enumerate(my_list):   
        if idx % 2 == 0 and value == '':
            my_list[idx] = '\n'
            my_list.insert(idx + 1, '\n')

    for idx, value in enumerate(my_list):   

        if idx % 2 == 0:
            if value == '':
                slide_list.append('\n')
                slide_list.append('\n')
                idx_count += 1
            else:
                slide_list.append(value)
        else:
            if value == '\n':
                slide_list[idx_count] = slide_list[idx_count] + value
                idx_count += 1
            else:
                slide_list[idx_count] = slide_list[idx_count] + "\n" + value
                idx_count += 1

    return slide_list
    
    